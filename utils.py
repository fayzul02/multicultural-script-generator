import os
from pptx import Presentation
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_core.prompts import PromptTemplate, ChatPromptTemplate
from langchain_core.output_parsers import StrOutputParser

# 1. Fungsi Parsing PPT
def extract_text_from_ppt(uploaded_file):
    prs = Presentation(uploaded_file)
    slides_content = []
    for i, slide in enumerate(prs.slides):
        text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + " "
        clean_text = text.strip() if text.strip() else "[Slide ini hanya gambar/grafik]"
        slides_content.append({"slide_number": i + 1, "content": clean_text})
    return slides_content

# 2. Fungsi Generasi Script (DIPERBAIKI)
def generate_cultural_script(slide_text, target_culture, api_key):
    if not api_key:
        return "⚠️ API Key belum dimasukkan."

    # Setup Model
    llm = ChatGoogleGenerativeAI(
        model="gemma-3-12b",
        google_api_key=api_key,
        temperature=0.7
    )

    # Setup Prompt
    template = """
    Kamu adalah seorang Ahli Komunikasi Bisnis Internasional dan Pakar Budaya {culture}.
    Tugasmu adalah membuat naskah presentasi (script) lisan untuk slide berikut.

    KONTEKS SLIDE:
    "{text}"

    INSTRUKSI KHUSUS:
    1. Buat naskah formal namun natural untuk diucapkan.
    2. ADAPTASI BUDAYA: Wajib masukkan elemen budaya {culture} (seperti idiom, peribahasa, atau gaya sapaan khas) yang relevan dengan topik bisnis di slide ini.
    3. Jelaskan singkat di akhir, kenapa kamu memilih adaptasi budaya tersebut.

    FORMAT OUTPUT:
    [Naskah]: (Tulis naskah di sini)
    [Catatan Budaya]: (Penjelasan adaptasi)
    """

    prompt = PromptTemplate(
        input_variables=["culture", "text"],
        template=template
    )

    # Chain Baru (LCEL)
    chain = prompt | llm | StrOutputParser()
    
    # --- BAGIAN INI YANG PENTING ---
    # Gunakan .invoke(), JANGAN .run()
    try:
        response = chain.invoke({"culture": target_culture, "text": slide_text})
        return response
    except Exception as e:
        return f"Terjadi kesalahan: {e}"